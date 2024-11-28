# utils/file_processing.py

import io
import pdfplumber
import docx
from pptx import Presentation

def extract_text_from_pdf(file_bytes):
    """Extracts text from a PDF file, including any tables."""
    text = ""
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n"

                # Extract tables as text
                tables = page.extract_tables()
                for table in tables:
                    table_text = "\n".join(["\t".join([str(cell) for cell in row]) for row in table])
                    text += table_text + "\n"
    except Exception as e:
        text += f"\nError extracting text from PDF: {e}\n"
    return text

def extract_text_from_docx(file_bytes):
    """Extracts text from a DOCX file, including any tables."""
    text = ""
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text for cell in row.cells)
                text += row_text + "\n"
    except Exception as e:
        text += f"\nError extracting text from DOCX: {e}\n"
    return text

def extract_text_from_pptx(file_bytes):
    """Extracts text from a PPTX file, including all slides and shapes."""
    text = ""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text += f"\nError extracting text from PPTX: {e}\n"
    return text

def extract_all_text(files):
    """Extracts text from a list of uploaded files, handling PDF, DOCX, and PPTX formats."""
    extracted_text = {}
    for file in files:
        # Check that file has 'name' and 'read' attributes (i.e., is an UploadedFile)
        if not hasattr(file, 'name') or not hasattr(file, 'read'):
            extracted_text["Unknown"] = "Error: Expected UploadedFile object, but got bytes."
            continue
        
        filename = file.name  # Define filename before the try block
        try:
            file_bytes = file.read()
            
            if filename.lower().endswith('.pdf'):
                extracted_text[filename] = extract_text_from_pdf(file_bytes)
            elif filename.lower().endswith('.docx'):
                extracted_text[filename] = extract_text_from_docx(file_bytes)
            elif filename.lower().endswith('.pptx'):
                extracted_text[filename] = extract_text_from_pptx(file_bytes)
            else:
                extracted_text[filename] = "Unsupported file format."
        
        except Exception as e:
            extracted_text[filename] = f"Error processing file {filename}: {e}"
    return extracted_text

